from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from institutional_bp_engine import build_projection, clean_intake

NAVY = RGBColor(6, 26, 58)
BLUE = RGBColor(11, 95, 255)
GREY = RGBColor(100, 116, 139)
LIGHT = RGBColor(247, 250, 252)
LINE = RGBColor(226, 232, 240)

def _title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    box.text_frame.text = title
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(12), Inches(0.35))
        st.text_frame.text = subtitle
        st.text_frame.paragraphs[0].font.size = Pt(10)
        st.text_frame.paragraphs[0].font.color.rgb = GREY

def _footer(slide, page):
    f = slide.shapes.add_textbox(Inches(0.55), Inches(7.07), Inches(12.3), Inches(0.25))
    f.text_frame.text = f"Private & Confidential | MG Advisory | {page}"
    f.text_frame.paragraphs[0].font.size = Pt(8)
    f.text_frame.paragraphs[0].font.color.rgb = GREY

def _card(slide, x, y, w, h, title, body):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = LIGHT
    s.line.color.rgb = LINE
    t = slide.shapes.add_textbox(Inches(x+0.15), Inches(y+0.12), Inches(w-0.3), Inches(0.32))
    t.text_frame.text = str(title)
    t.text_frame.paragraphs[0].font.size = Pt(11)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = NAVY
    b = slide.shapes.add_textbox(Inches(x+0.15), Inches(y+0.50), Inches(w-0.3), Inches(h-0.55))
    b.text_frame.text = str(body)
    b.text_frame.paragraphs[0].font.size = Pt(9)
    b.text_frame.paragraphs[0].font.color.rgb = GREY

def build_institutional_im_deck(project: Dict[str, Any], intake: Dict[str, Any], output_path: str) -> str:
    x = clean_intake(intake)
    rows = build_projection(x)["projection"]
    latest = rows[0] if rows else {}
    final = rows[-1] if rows else latest

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = []

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    tx = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(10.0), Inches(1.0))
    tx.text_frame.text = x["company_name"]
    tx.text_frame.paragraphs[0].font.size = Pt(42)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(9), Inches(0.6))
    sub.text_frame.text = f"{x['deal_type']} Investment Memorandum | {x['sector']}"
    sub.text_frame.paragraphs[0].font.size = Pt(19)
    sub.text_frame.paragraphs[0].font.color.rgb = BLUE
    _footer(slide, 1)

    # 2 Contents
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Contents", "Institutional IM structure")
    sections = ["Executive Summary", "Investment Highlights", "Business Overview", "Market & Customers", "Operations", "Financial Overview", "Value Creation", "QoE Focus", "Risks", "Process & Next Steps"]
    for i, sec in enumerate(sections):
        _card(slide, 0.75 + (i % 2)*6.05, 1.25 + (i//2)*1.05, 5.65, 0.75, f"{i+1:02d}", sec)
    _footer(slide, 2)

    # 3 Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Executive Summary", "Key facts and investment context")
    _card(slide, 0.75, 1.25, 2.8, 1.0, "Revenue", f"{x['currency']} {latest.get('revenue', 0):,.1f}")
    _card(slide, 3.85, 1.25, 2.8, 1.0, "EBITDA", f"{x['currency']} {latest.get('ebitda', 0):,.1f}")
    _card(slide, 6.95, 1.25, 2.8, 1.0, "Final EBITDA", f"{x['currency']} {final.get('ebitda', 0):,.1f}")
    _card(slide, 10.05, 1.25, 2.4, 1.0, "FTE", f"{x.get('fte', 'N/A')}")
    _card(slide, 0.75, 2.65, 11.7, 1.1, "Investment thesis", x.get("investment_thesis") or "To be completed during management session.")
    _card(slide, 0.75, 4.05, 5.65, 1.25, "Why now", "Professionalised reporting, BP, QoE and investor documentation accelerate diligence readiness.")
    _card(slide, 6.8, 4.05, 5.65, 1.25, "Diligence focus", "Revenue quality, margin sustainability, working capital, capex, leverage and execution risk.")
    _footer(slide, 3)

    # 4 Investment Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Investment Highlights", "Core investment messages")
    highlights = [
        ("Scale", f"{x['currency']} {latest.get('revenue',0):,.1f} revenue platform in {x['sector']}"),
        ("Margin", f"{latest.get('ebitda_margin',0):.1%} EBITDA margin with operating leverage potential"),
        ("Cash", f"{x['currency']} {latest.get('fcf',0):,.1f} FCF in first forecast year"),
        ("Value Creation", x.get("value_creation_plan") or "Revenue growth, margin improvement, NWC discipline and capex prioritisation"),
        ("Risk", x.get("key_risks") or "Forecast achievability, customer concentration and cash conversion to validate")
    ]
    for i, (a,b) in enumerate(highlights):
        _card(slide, 0.75 + (i%2)*6.05, 1.25 + (i//2)*1.35, 5.65, 1.05, a, b)
    _footer(slide, 4)

    # 5 Revenue Build
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Revenue Build-Up", "Volume, price, contracted revenue and new business")
    table = slide.shapes.add_table(len(rows)+1, 5, Inches(0.75), Inches(1.4), Inches(11.8), Inches(3.5)).table
    headers = ["Year", "Revenue", "Contracted", "New Business", "Churn Impact"]
    for c,h in enumerate(headers): table.cell(0,c).text = h
    for r,row in enumerate(rows, start=1):
        vals = [row["year"], row["revenue"], row["contracted_revenue"], row["new_business"], row["churn_impact"]]
        for c,v in enumerate(vals): table.cell(r,c).text = f"{v:,.1f}" if isinstance(v, (int,float)) else str(v)
    _card(slide, 0.75, 5.35, 11.8, 0.8, "Commercial diligence questions", "Validate customer contracts, renewal terms, pipeline conversion, pricing and volume assumptions.")
    _footer(slide, 5)

    # 6 Operations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Operations & Capacity", "Utilisation, FTE and operating leverage")
    _card(slide, 0.75, 1.35, 2.8, 1.0, "Capacity", f"{x.get('capacity',0):,.0f}")
    _card(slide, 3.85, 1.35, 2.8, 1.0, "Utilisation", f"{x.get('utilisation',0):.1%}")
    _card(slide, 6.95, 1.35, 2.8, 1.0, "FTE", f"{x.get('fte',0):,.0f}")
    _card(slide, 10.05, 1.35, 2.4, 1.0, "Rev/FTE", f"{latest.get('revenue_per_fte',0):,.1f}")
    _card(slide, 0.75, 2.85, 11.7, 1.2, "Operating leverage", "Incremental volumes are expected to improve fixed-cost absorption and EBITDA conversion, subject to validation of capacity, bottlenecks and capex.")
    _footer(slide, 6)

    # 7 Financial Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Financial Overview", "Projected P&L and cash generation")
    table = slide.shapes.add_table(len(rows)+1, 7, Inches(0.55), Inches(1.35), Inches(12.25), Inches(3.7)).table
    headers = ["Year", "Revenue", "Gross Profit", "EBITDA", "EBITDA %", "FCF", "Net Debt/EBITDA"]
    for c,h in enumerate(headers): table.cell(0,c).text = h
    for r,row in enumerate(rows, start=1):
        vals = [row["year"], row["revenue"], row["gross_profit"], row["ebitda"], row["ebitda_margin"], row["fcf"], row["net_debt_ebitda"]]
        for c,v in enumerate(vals):
            table.cell(r,c).text = f"{v:.1%}" if c==4 and isinstance(v,(int,float)) else (f"{v:,.1f}" if isinstance(v,(int,float)) else str(v))
    _footer(slide, 7)

    # 8 Value Creation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Value Creation Plan", "Commercial, operational and financial levers")
    _card(slide, 0.75, 1.35, 11.7, 0.95, "Management plan", x.get("value_creation_plan") or "To be completed during management workshop.")
    _card(slide, 0.75, 2.65, 3.6, 1.4, "Commercial", "Growth, price, volume, mix, contract conversion, churn reduction.")
    _card(slide, 4.85, 2.65, 3.6, 1.4, "Operational", "Utilisation, procurement, productivity, energy, logistics.")
    _card(slide, 8.95, 2.65, 3.5, 1.4, "Financial", "NWC, capex discipline, debt optimisation and cash conversion.")
    _footer(slide, 8)

    # 9 QoE Focus
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Quality of Earnings Focus", "Areas to validate during diligence")
    qoe = ["Revenue recognition and cut-off", "Customer concentration and retention", "One-off costs and run-rate adjustments", "Working capital normalisation", "Debt-like items and off-balance sheet liabilities"]
    for i,item in enumerate(qoe):
        _card(slide, 0.75, 1.25 + i*0.9, 11.7, 0.7, f"QoE area {i+1}", item)
    _footer(slide, 9)

    # 10 Risks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Risks & Mitigants", "Investment committee diligence framing")
    risks = x.get("key_risks") or "Forecast delivery, customer concentration, margin pressure, capex needs, liquidity and leverage capacity."
    _card(slide, 0.75, 1.35, 11.7, 1.2, "Key risks", risks)
    _card(slide, 0.75, 2.95, 5.65, 1.35, "Mitigants", "Independent QoE, customer contract review, monthly trading review, downside case and covenant sensitivity.")
    _card(slide, 6.8, 2.95, 5.65, 1.35, "Open diligence questions", x.get("diligence_questions") or "To be completed with advisor and management.")
    _footer(slide, 10)

    # 11 Process
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "Process & Next Steps", "Recommended workplan")
    steps = ["Complete data room upload", "Validate BP assumptions", "Run QoE and debt-like item review", "Populate customer / market sections", "Finalise IC materials and management Q&A"]
    for i,s in enumerate(steps):
        _card(slide, 0.75, 1.25 + i*0.9, 11.7, 0.7, f"Step {i+1}", s)
    _footer(slide, 11)

    prs.save(output_path)
    return output_path
