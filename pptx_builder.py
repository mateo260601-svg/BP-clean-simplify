"""pptx_builder.py — editable institutional PowerPoint IM generator."""
from __future__ import annotations
from typing import Dict, List, Any, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

NAVY=RGBColor(13,27,62); BLUE=RGBColor(0,61,255); CYAN=RGBColor(0,174,239); LIGHT=RGBColor(244,247,251); GREY=RGBColor(100,116,139); WHITE=RGBColor(255,255,255); BLACK=RGBColor(15,23,42); RED=RGBColor(185,28,28); GREEN=RGBColor(5,150,105)

def _blank(prs): return prs.slide_layouts[6]
def _tx(shape, text, size=10, bold=False, color=BLACK, align=None):
    tf=shape.text_frame; tf.clear(); p=tf.paragraphs[0]
    if align is not None: p.alignment=align
    r=p.add_run(); r.text=str(text); r.font.name="Arial"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return tf

def _add_title(slide, title, subtitle=None):
    box=slide.shapes.add_textbox(Inches(.55), Inches(.35), Inches(12.2), Inches(.55)); _tx(box, title, 24, True, NAVY)
    line=slide.shapes.add_shape(1, Inches(.55), Inches(1.02), Inches(12.2), Inches(.03)); line.fill.solid(); line.fill.fore_color.rgb=BLUE; line.line.color.rgb=BLUE
    if subtitle:
        sub=slide.shapes.add_textbox(Inches(.58), Inches(1.08), Inches(11.8), Inches(.3)); _tx(sub, subtitle, 9, False, BLUE)

def _footer(slide, page):
    box=slide.shapes.add_textbox(Inches(.55), Inches(7.16), Inches(6.2), Inches(.22)); _tx(box, "Private & Confidential – Not For Distribution", 7, False, GREY)
    pg=slide.shapes.add_textbox(Inches(12.25), Inches(7.12), Inches(.45), Inches(.25)); _tx(pg, str(page), 8, True, BLUE, PP_ALIGN.RIGHT)

def _fmt_num(v, cur="€"):
    try: v=float(v)
    except Exception: return "—"
    sign="-" if v<0 else ""; v=abs(v)
    if v>=1000: return f"{sign}{cur}{v/1000:,.1f}m"
    return f"{sign}{cur}{v:,.0f}k"

def _fmt_pct(v):
    try: return f"{float(v)*100:.1f}%"
    except Exception: return "—"

def _metric(fin,k,p):
    return (fin.get(k) or {}).get(str(p)) if p is not None else None

def build_im_deck(model: Dict[str,Any], slides: List[Dict[str,Any]], output_path: str, template_path: Optional[str]=None) -> str:
    prs = Presentation(template_path) if template_path and Path(template_path).exists() else Presentation()
    prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    # Start from a clean file even if a template is supplied; template can be used later for masters.
    if len(prs.slides)>0:
        prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    page=1
    for s in slides:
        fn=globals().get(f"_slide_{s['id']}")
        if fn:
            if s['id'] in ('contents',): fn(prs, slides, page)
            else: fn(prs, model, page)
            page += 1
    prs.save(output_path); return output_path

def _slide_cover(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs))
    bg=slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)); bg.fill.solid(); bg.fill.fore_color.rgb=WHITE; bg.line.color.rgb=WHITE
    accent=slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.7), Inches(7.5)); accent.fill.solid(); accent.fill.fore_color.rgb=BLUE; accent.line.color.rgb=BLUE
    box=slide.shapes.add_textbox(Inches(.75), Inches(1.9), Inches(11), Inches(2.7)); tf=_tx(box, model.get('company_name','Target Company'), 38, True, WHITE); 
    p=tf.add_paragraph(); r=p.add_run(); r.text="Information Memorandum"; r.font.name="Arial"; r.font.size=Pt(24); r.font.color.rgb=WHITE
    p=tf.add_paragraph(); r=p.add_run(); r.text="Prepared from uploaded financial information"; r.font.name="Arial"; r.font.size=Pt(12); r.font.color.rgb=WHITE
    # Right KPI cards
    fin=model.get('financials',{}); periods=model.get('periods',[]); latest=periods[-1] if periods else None; cur=model.get('currency','€')
    cards=[('Revenue',_fmt_num(_metric(fin,'revenue',latest),cur)),('EBITDA',_fmt_num(_metric(fin,'ebitda',latest),cur)),('Data quality',(model.get('checks') or {}).get('quality_label','—'))]
    y=2.0
    for h,v in cards:
        rect=slide.shapes.add_shape(1, Inches(8.1), Inches(y), Inches(3.6), Inches(.82)); rect.fill.solid(); rect.fill.fore_color.rgb=LIGHT; rect.line.color.rgb=RGBColor(220,226,235)
        t=slide.shapes.add_textbox(Inches(8.32), Inches(y+.11), Inches(3.1), Inches(.5)); tf=_tx(t, v, 18, True, NAVY); p=tf.add_paragraph(); r=p.add_run(); r.text=h; r.font.name="Arial"; r.font.size=Pt(8); r.font.color.rgb=GREY
        y += 1.05
    _footer(slide,page)

def _slide_disclaimer(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Disclaimer")
    txt=("This presentation is confidential and has been prepared solely for discussion purposes. It does not constitute an offer, commitment, recommendation or investment advice. "
         "The analysis is based on information uploaded by the user and may include unaudited, incomplete or management-provided data. Recipients should conduct their own due diligence and verify all statements independently. "
         "Forward-looking information is subject to risks, uncertainties and assumptions, and actual results may differ materially.")
    box=slide.shapes.add_textbox(Inches(.75), Inches(1.45), Inches(11.9), Inches(4.8)); tf=_tx(box, txt, 12, False, BLACK); tf.word_wrap=True
    _footer(slide,page)

def _slide_contents(prs, slides, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Contents")
    y=1.45
    for i,s in enumerate(slides,1):
        circ=slide.shapes.add_shape(1, Inches(.8), Inches(y), Inches(.42), Inches(.35)); circ.fill.solid(); circ.fill.fore_color.rgb=BLUE; circ.line.color.rgb=BLUE
        _tx(circ, str(i), 10, True, WHITE, PP_ALIGN.CENTER)
        box=slide.shapes.add_textbox(Inches(1.45), Inches(y-.02), Inches(9.8), Inches(.35)); _tx(box, s['title'], 14, True, NAVY)
        y += .52
    _footer(slide,page)

def _slide_executive_summary(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Executive Summary", "Auto-generated transaction overview from uploaded accounts")
    bullets=_highlights(model)[:4]; x=.75
    for i,(h,b) in enumerate(bullets):
        y=1.5 + (i%2)*2.25; x=.75 + (i//2)*6.15
        rect=slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.55), Inches(1.75)); rect.fill.solid(); rect.fill.fore_color.rgb=LIGHT; rect.line.color.rgb=RGBColor(220,226,235)
        bar=slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(.12), Inches(1.75)); bar.fill.solid(); bar.fill.fore_color.rgb=BLUE; bar.line.color.rgb=BLUE
        title=slide.shapes.add_textbox(Inches(x+.28), Inches(y+.18), Inches(5.0), Inches(.3)); _tx(title,h,12,True,BLUE)
        body=slide.shapes.add_textbox(Inches(x+.28), Inches(y+.55), Inches(5.0), Inches(.9)); body.text_frame.word_wrap=True; _tx(body,b,9,False,BLACK)
    _footer(slide,page)

def _slide_investment_highlights(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Investment Highlights", "Key investment messages generated from standardised financials")
    y=1.35
    for i,(h,b) in enumerate(_highlights(model)[:5],1):
        icon=slide.shapes.add_shape(1, Inches(.75), Inches(y), Inches(.55), Inches(.55)); icon.fill.solid(); icon.fill.fore_color.rgb=BLUE; icon.line.color.rgb=BLUE; _tx(icon,str(i),14,True,WHITE,PP_ALIGN.CENTER)
        title=slide.shapes.add_textbox(Inches(1.55), Inches(y-.02), Inches(10.9), Inches(.28)); _tx(title,h,13,True,BLUE)
        body=slide.shapes.add_textbox(Inches(1.55), Inches(y+.32), Inches(10.8), Inches(.48)); body.text_frame.word_wrap=True; _tx(body,b,8.8,False,BLACK)
        y += 1.02
    _footer(slide,page)

def _highlights(model):
    fin=model.get('financials',{}); periods=model.get('periods',[]); latest=periods[-1] if periods else None; cur=model.get('currency','€')
    rev=_metric(fin,'revenue',latest); e=_metric(fin,'ebitda',latest); margin=(model.get('kpis',{}).get('ebitda_margin',{}) or {}).get(str(latest))
    return [
      ("Standardised financial history", "Audited accounts and/or management accounts are mapped into a consistent P&L, balance sheet, cash and working-capital framework."),
      ("Revenue base identified", f"Latest available revenue is {_fmt_num(rev,cur)}, with historical trend and growth metrics prepared for diligence."),
      ("Profitability profile analysed", f"Latest available EBITDA is {_fmt_num(e,cur)} with EBITDA margin of {_fmt_pct(margin)} where data is available."),
      ("Liquidity and leverage view", "Cash, debt and net debt metrics are analysed when balance sheet information is present in the uploaded files."),
      ("Diligence-ready outputs", "The generated PowerPoint and Excel outputs are editable, institutional and structured for investor / lender discussions."),
    ]

def _slide_financial_overview(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Financial Overview", "Revenue and EBITDA development")
    fin=model.get('financials',{}); periods=model.get('periods',[]); cur=model.get('currency','€')
    data=CategoryChartData(); data.categories=periods
    data.add_series('Revenue',[float(_metric(fin,'revenue',p) or 0) for p in periods])
    data.add_series('EBITDA',[float(_metric(fin,'ebitda',p) or 0) for p in periods])
    chart=slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(.75), Inches(1.45), Inches(7.35), Inches(4.45), data).chart
    chart.has_legend=True; chart.legend.position=XL_LEGEND_POSITION.BOTTOM
    latest=periods[-1] if periods else None
    cards=[('Revenue',_fmt_num(_metric(fin,'revenue',latest),cur)),('EBITDA',_fmt_num(_metric(fin,'ebitda',latest),cur)),('Revenue Growth',_fmt_pct((model.get('kpis',{}).get('revenue_growth',{}) or {}).get(str(latest)))),('EBITDA Margin',_fmt_pct((model.get('kpis',{}).get('ebitda_margin',{}) or {}).get(str(latest))))]
    y=1.45
    for h,v in cards:
        rect=slide.shapes.add_shape(1, Inches(8.55), Inches(y), Inches(3.8), Inches(.82)); rect.fill.solid(); rect.fill.fore_color.rgb=LIGHT; rect.line.color.rgb=RGBColor(220,226,235)
        tf=_tx(slide.shapes.add_textbox(Inches(8.76), Inches(y+.11), Inches(3.3), Inches(.55)), v, 16, True, NAVY); p=tf.add_paragraph(); r=p.add_run(); r.text=h; r.font.name="Arial"; r.font.size=Pt(8); r.font.color.rgb=GREY
        y += 1.0
    _footer(slide,page)

def _slide_profitability(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Profitability & Margin Profile")
    fin=model.get('financials',{}); periods=model.get('periods',[]); k=model.get('kpis',{})
    rows=[('Revenue','revenue',None),('Gross Profit','gross_profit','gross_margin'),('EBITDA','ebitda','ebitda_margin'),('EBIT','ebit','ebit_margin'),('Net Income','net_income','net_margin')]
    _table(slide, rows, periods, fin, k, Inches(.75), Inches(1.45), Inches(11.9), model.get('currency','€'))
    _footer(slide,page)

def _table(slide, rows, periods, fin, kpis, x, y, w, cur):
    cols=1+len(periods); table=slide.shapes.add_table(len(rows)+1, cols, x, y, w, Inches(.42*(len(rows)+1))).table
    table.cell(0,0).text="Metric"
    for j,p in enumerate(periods,1): table.cell(0,j).text=str(p)
    for j in range(cols):
        c=table.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY; _tx(c.text_frame.paragraphs[0]._parent, c.text, 8, True, WHITE) if False else None
        for p in c.text_frame.paragraphs:
            for r in p.runs: r.font.color.rgb=WHITE; r.font.bold=True; r.font.size=Pt(8)
    for i,(label,key,kpi) in enumerate(rows,1):
        table.cell(i,0).text=label
        for j,p in enumerate(periods,1):
            v=(kpis.get(kpi,{}) if kpi else fin.get(key,{})).get(str(p)) if (kpi or key) else None
            table.cell(i,j).text=_fmt_pct(v) if kpi else _fmt_num(v,cur)
    for row in table.rows:
        for c in row.cells:
            c.margin_left=Inches(.06); c.margin_right=Inches(.06)
            for p in c.text_frame.paragraphs:
                for r in p.runs: r.font.name="Arial"; r.font.size=Pt(8)

def _slide_cash_debt_liquidity(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Cash, Debt & Liquidity", "Net debt and leverage overview")
    fin=model.get('financials',{}); periods=model.get('periods',[]); cur=model.get('currency','€'); k=model.get('kpis',{})
    rows=[('Cash','cash',None),('Gross Debt','debt',None),('Net Debt','net_debt','net_debt'),('Net Debt / EBITDA','net_debt_ebitda','net_debt_ebitda')]
    # custom table
    table=slide.shapes.add_table(len(rows)+1, 1+len(periods), Inches(.75), Inches(1.45), Inches(11.9), Inches(2.5)).table
    table.cell(0,0).text="Metric"
    for j,p in enumerate(periods,1): table.cell(0,j).text=str(p)
    for i,(label,key,kpi) in enumerate(rows,1):
        table.cell(i,0).text=label
        for j,p in enumerate(periods,1):
            if kpi == 'net_debt_ebitda':
                val=(k.get('net_debt_ebitda',{}) or {}).get(str(p)); table.cell(i,j).text='—' if val is None else f"{val:.1f}x"
            elif kpi == 'net_debt': table.cell(i,j).text=_fmt_num((k.get('net_debt',{}) or {}).get(str(p)),cur)
            else: table.cell(i,j).text=_fmt_num((fin.get(key,{}) or {}).get(str(p)),cur)
    _style_table(table)
    _footer(slide,page)

def _style_table(table):
    for i,row in enumerate(table.rows):
        for cell in row.cells:
            if i==0: cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name="Arial"; r.font.size=Pt(8); r.font.bold=(i==0); r.font.color.rgb=WHITE if i==0 else BLACK

def _slide_working_capital(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Working Capital", "Inventory, receivables and payables development")
    fin=model.get('financials',{}); periods=model.get('periods',[]); cur=model.get('currency','€')
    rows=[('Inventory','inventory',None),('Receivables','receivables',None),('Payables','payables',None),('Net Working Capital','nwc','nwc'),('NWC % Sales','nwc_pct_sales','nwc_pct_sales')]
    table=slide.shapes.add_table(len(rows)+1,1+len(periods),Inches(.75),Inches(1.45),Inches(11.9),Inches(2.8)).table
    table.cell(0,0).text='Metric'
    for j,p in enumerate(periods,1): table.cell(0,j).text=str(p)
    k=model.get('kpis',{})
    for i,(label,key,kpi) in enumerate(rows,1):
        table.cell(i,0).text=label
        for j,p in enumerate(periods,1):
            val=(k.get(kpi,{}) if kpi else fin.get(key,{})).get(str(p))
            table.cell(i,j).text=_fmt_pct(val) if kpi=='nwc_pct_sales' else _fmt_num(val,cur)
    _style_table(table); _footer(slide,page)

def _slide_value_creation(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Value Creation Opportunity", "Illustrative levers derived from financial profile")
    levers=[('Revenue growth','Volume, price and mix improvement'),('Margin expansion','COGS discipline and operating leverage'),('Working capital release','Receivables, inventory and payment terms optimisation'),('Cash generation','Capex discipline and EBITDA conversion')]
    x=.75
    for i,(h,b) in enumerate(levers):
        rect=slide.shapes.add_shape(1, Inches(x+i*3.05), Inches(1.6), Inches(2.65), Inches(3.7)); rect.fill.solid(); rect.fill.fore_color.rgb=LIGHT if i else BLUE; rect.line.color.rgb=BLUE
        _tx(slide.shapes.add_textbox(Inches(x+i*3.05+.18), Inches(1.95), Inches(2.25), Inches(.45)), h, 14, True, WHITE if i==0 else NAVY)
        body=slide.shapes.add_textbox(Inches(x+i*3.05+.18), Inches(2.55), Inches(2.25), Inches(1.3)); body.text_frame.word_wrap=True; _tx(body,b,9,False,WHITE if i==0 else BLACK)
    _footer(slide,page)

def _slide_data_quality(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Data Quality & Diligence Items")
    checks=model.get('checks',{}) or {}; available=checks.get('available_metrics',[]); missing=checks.get('missing_critical_metrics',[])
    texts=[('Quality score', f"{checks.get('quality_score','—')} / 100"),('Available metrics', ', '.join(available) if available else 'None identified'),('Missing critical metrics', ', '.join(missing) if missing else 'None'),('Recommended next step','Upload complete audited accounts, latest management accounts, debt schedule, cash report and customer/operational KPIs.')]
    y=1.4
    for h,b in texts:
        _tx(slide.shapes.add_textbox(Inches(.8), Inches(y), Inches(2.7), Inches(.3)), h, 10, True, BLUE)
        box=slide.shapes.add_textbox(Inches(3.25), Inches(y), Inches(9.3), Inches(.55)); box.text_frame.word_wrap=True; _tx(box,b,10,False,BLACK)
        y += .8
    _footer(slide,page)

def _slide_process_next_steps(prs, model, page):
    slide=prs.slides.add_slide(_blank(prs)); _add_title(slide,"Process & Next Steps")
    steps=[('Week 1','Document upload & validation'),('Week 2','Financial model and QoE review'),('Week 3','Draft IM / lender materials'),('Week 4','Management review and revisions'),('Week 5+','Investor outreach / diligence')]
    x=.8
    for i,(h,b) in enumerate(steps):
        circ=slide.shapes.add_shape(1, Inches(x+i*2.45), Inches(2.0), Inches(.65), Inches(.65)); circ.fill.solid(); circ.fill.fore_color.rgb=BLUE; circ.line.color.rgb=BLUE; _tx(circ,str(i+1),15,True,WHITE,PP_ALIGN.CENTER)
        _tx(slide.shapes.add_textbox(Inches(x+i*2.45-.15), Inches(2.85), Inches(1.4), Inches(.3)), h, 11, True, NAVY, PP_ALIGN.CENTER)
        body=slide.shapes.add_textbox(Inches(x+i*2.45-.3), Inches(3.25), Inches(1.8), Inches(.8)); body.text_frame.word_wrap=True; _tx(body,b,8,False,BLACK,PP_ALIGN.CENTER)
    _footer(slide,page)
